#!/usr/bin/env python3
"""
Script para generar presentación PowerPoint para venta a HomeBanking
Enfoque: Por qué un banco debería integrar seguros con AI en su plataforma
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Colores corporativos
DARK_BLUE = RGBColor(26, 54, 93)      # #1a365d
MEDIUM_BLUE = RGBColor(43, 108, 176)  # #2b6cb0
LIGHT_BLUE = RGBColor(190, 227, 248)  # #bee3f8
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(74, 85, 104)     # #4a5568
GREEN = RGBColor(56, 161, 105)        # #38a169
ORANGE = RGBColor(237, 137, 54)       # #ed8936

def add_title_slide(prs, title, subtitle):
    """Agregar slide de título"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """Agregar slide de sección"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = MEDIUM_BLUE
    background.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, has_highlight=False, highlight_text=""):
    """Agregar slide de contenido con bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
        p.level = 0
    
    # Highlight box si existe
    if has_highlight and highlight_text:
        highlight_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(9), Inches(0.8)
        )
        highlight_box.fill.solid()
        highlight_box.fill.fore_color.rgb = LIGHT_BLUE
        highlight_box.line.fill.background()
        
        hl_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(8.6), Inches(0.5))
        tf = hl_text.text_frame
        p = tf.paragraphs[0]
        p.text = highlight_text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_stats_slide(prs, title, stats):
    """Agregar slide con estadísticas"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Estadísticas en grid
    num_stats = len(stats)
    cols = 2 if num_stats <= 4 else 3
    width = Inches(4) if cols == 2 else Inches(3)
    
    for i, (value, label, source) in enumerate(stats):
        row = i // cols
        col = i % cols
        x = Inches(0.5) + col * (width + Inches(0.3))
        y = Inches(1.8) + row * Inches(2.2)
        
        # Caja de estadística
        stat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, Inches(1.8)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = LIGHT_BLUE
        stat_box.line.fill.background()
        
        # Valor
        val_box = slide.shapes.add_textbox(x, y + Inches(0.2), width, Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = MEDIUM_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.9), width - Inches(0.2), Inches(0.5))
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        # Source
        src_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.4), width - Inches(0.2), Inches(0.3))
        tf = src_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Fuente: {source}"
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Agregar slide con dos columnas"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Columna izquierda - título
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = MEDIUM_BLUE
    
    # Columna izquierda - contenido
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4.5), Inches(4))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Línea separadora
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(1.5), Inches(0.02), Inches(4.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()
    
    # Columna derecha - título
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.5), Inches(4.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    # Columna derecha - contenido
    right_content = slide.shapes.add_textbox(Inches(5.3), Inches(2.1), Inches(4.5), Inches(4))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    return slide

def add_product_slide(prs, title, products):
    """Agregar slide de productos/herramientas"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Productos en grid
    for i, (name, description, icon) in enumerate(products):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * Inches(4.8)
        y = Inches(1.5) + row * Inches(1.7)
        
        # Caja de producto
        prod_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1.5)
        )
        prod_box.fill.solid()
        prod_box.fill.fore_color.rgb = WHITE
        prod_box.line.color.rgb = MEDIUM_BLUE
        prod_box.line.width = Pt(2)
        
        # Icono/Número
        icon_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5)
        )
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = MEDIUM_BLUE
        icon_box.line.fill.background()
        
        icon_text = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.22), Inches(0.5), Inches(0.4))
        tf = icon_text.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Nombre
        name_box = slide.shapes.add_textbox(x + Inches(0.75), y + Inches(0.15), Inches(3.5), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        # Descripción
        desc_box = slide.shapes.add_textbox(x + Inches(0.75), y + Inches(0.6), Inches(3.5), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
    
    return slide

def create_presentation():
    """Crear la presentación completa"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Portada
    add_title_slide(
        prs,
        "Seguros Inteligentes\npara HomeBanking",
        "Cómo integrar AI para aumentar ingresos\ny mejorar la experiencia del cliente"
    )
    
    # Slide 2: La Oportunidad
    add_section_slide(prs, "La Oportunidad")
    
    # Slide 3: El Mercado
    add_stats_slide(prs, "El Mercado de Seguros está Transformándose", [
        ("88%", "de aseguradoras ya usan AI en al menos una función", "McKinsey 2025"),
        ("58%", "de clientes quieren consejos personalizados", "Accenture 2026"),
        ("$5.2T", "mercado global de seguros", "Swiss Re"),
        ("14%→70%", "adopción de AI en underwriting en 3 años", "Accenture"),
    ])
    
    # Slide 4: Por qué un Banco
    add_content_slide(
        prs,
        "¿Por qué un Banco debería ofrecer Seguros?",
        [
            "Ya tienen la confianza del cliente (relación establecida)",
            "Datos financieros que mejoran el risk scoring",
            "Canal digital establecido (app, web) con millones de usuarios",
            "Oportunidad de cross-sell en momentos clave (créditos, hipotecas)",
            "Nuevo revenue stream sin inversión en infraestructura de seguros"
        ],
        has_highlight=True,
        highlight_text="El banco se convierte en distribuidor inteligente, no en aseguradora"
    )
    
    # Slide 5: Problema Actual
    add_two_column_slide(
        prs,
        "Situación Actual vs. Propuesta",
        "HOY: Modelo Tradicional",
        [
            "Cliente debe salir del app para cotizar",
            "Formularios largos y tediosos",
            "Sin personalización",
            "Agentes no conocen al cliente",
            "Proceso toma días o semanas",
            "Baja conversión"
        ],
        "MAÑANA: Con AI Integrada",
        [
            "Cotización en el mismo app bancario",
            "Conversación natural en 3 minutos",
            "Recomendaciones basadas en perfil",
            "Contexto financiero del cliente",
            "Emisión inmediata",
            "Alta conversión"
        ]
    )
    
    # Slide 6: Nuestra Solución
    add_section_slide(prs, "Nuestra Solución")
    
    # Slide 7: Suite de Productos
    add_product_slide(prs, "Suite de Herramientas de AI para Seguros", [
        ("Motor de Cotización Conversacional", "Chat inteligente que recopila datos y genera cotizaciones personalizadas en minutos", "1"),
        ("Copiloto para Agentes", "Asistente AI que ayuda a agentes con contexto del cliente y recomendaciones", "2"),
        ("Analizador de Productos", "Explica productos complejos y compara opciones para el cliente", "3"),
        ("Agente de Claims", "Procesa siniestros automáticamente: triaje, fotos, pagos rápidos", "4"),
        ("Document Intelligence", "Extrae datos de documentos (DNI, facturas, fotos) automáticamente", "5"),
        ("Dashboard de Gobernanza", "Monitoreo, auditoría y compliance de decisiones de AI", "6"),
    ])
    
    # Slide 8: Integración Simple
    add_content_slide(
        prs,
        "Integración Simple con su HomeBanking",
        [
            "API REST documentada para integración en días, no meses",
            "Widget embebible en su app móvil y web",
            "White-label: su marca, nuestra tecnología",
            "Compatible con múltiples aseguradoras (multi-carrier)",
            "Cumplimiento regulatorio incluido",
            "Sin impacto en sus sistemas core"
        ],
        has_highlight=True,
        highlight_text="Time-to-market: 4-6 semanas"
    )
    
    # Slide 9: Resultados Probados
    add_stats_slide(prs, "Resultados Probados en Producción", [
        ("91%", "de intake de claims automatizado", "Dearborn Labs"),
        ("93%", "de pólizas emitidas sin intervención humana", "Dearborn Labs"),
        ("3X", "más eficiencia en manejo de claims", "Dearborn Labs"),
        (">30%", "ahorro en costos de underwriting", "BCG"),
    ])
    
    # Slide 10: Modelo de Negocio
    add_content_slide(
        prs,
        "Modelo de Negocio: Gana-Gana",
        [
            "Sin inversión inicial: setup incluido",
            "Revenue share por póliza vendida a través del canal",
            "El banco obtiene comisión por cada venta",
            "Métricas transparentes en dashboard en tiempo real",
            "Escalable: más productos, más ingresos"
        ]
    )
    
    # Slide 11: ROI
    add_two_column_slide(
        prs,
        "Proyección de Impacto Financiero",
        "Reducción de Costos",
        [
            "Automatización de cotización: -80% costo por quote",
            "Claims rápidos: menos reservas",
            "Sin call center para seguros básicos",
            "Menos errores de data entry"
        ],
        "Aumento de Ingresos",
        [
            "Nuevo revenue stream (comisiones)",
            "Cross-sell en momentos clave",
            "Mayor retención de clientes",
            "Upsell inteligente basado en datos",
            "NPS mejorado = más referidos"
        ]
    )
    
    # Slide 12: Casos de Uso
    add_content_slide(
        prs,
        "Casos de Uso en HomeBanking",
        [
            "Al aprobar un crédito automotriz → Ofrecer seguro de auto",
            "Al firmar hipoteca → Ofrecer seguro de hogar + vida",
            "Al detectar viaje internacional → Ofrecer seguro de viaje",
            "Al cumplir 40 años → Ofrecer seguro de vida/retiro",
            "Al reportar robo de tarjeta → Activar claim de seguro digital"
        ],
        has_highlight=True,
        highlight_text="Seguros contextuales = mayor conversión"
    )
    
    # Slide 13: Diferenciadores
    add_content_slide(
        prs,
        "¿Por qué Nosotros?",
        [
            "MVP funcionando hoy (no PowerPoint vaporware)",
            "Multi-LLM: OpenAI, Anthropic, modelos locales",
            "Arquitectura modular: use solo lo que necesita",
            "Agnóstico de carrier: integramos con quien usted elija",
            "Equipo con experiencia en seguros Y en AI",
            "Enfoque en compliance desde el diseño"
        ]
    )
    
    # Slide 14: Seguridad
    add_content_slide(
        prs,
        "Seguridad y Compliance",
        [
            "Datos encriptados en tránsito y en reposo",
            "No almacenamos datos sensibles del cliente",
            "Auditoría completa de decisiones de AI",
            "Preparado para regulaciones de AI (EU AI Act, IMDA)",
            "SOC2 Type II en roadmap",
            "Hosting flexible: cloud o on-premise"
        ]
    )
    
    # Slide 15: Roadmap
    add_content_slide(
        prs,
        "Roadmap de Implementación",
        [
            "Semana 1-2: Discovery y definición de productos",
            "Semana 3-4: Integración de APIs y configuración",
            "Semana 5: Testing con usuarios internos",
            "Semana 6: Piloto controlado (soft launch)",
            "Semana 8+: Rollout completo y optimización continua"
        ]
    )
    
    # Slide 16: CTA
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "¿Listo para transformar su\noferta de seguros?"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Agendemos una demo personalizada"
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Contacto
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "contacto@empresa.com"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Guardar
    output_path = Path(__file__).parent / "pitch-seguros-homebanking.pptx"
    prs.save(output_path)
    print(f"Presentación generada exitosamente: {output_path}")

if __name__ == "__main__":
    create_presentation()
