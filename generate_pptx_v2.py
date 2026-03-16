#!/usr/bin/env python3
"""
Presentación PowerPoint v2 - Pitch para HomeBanking Paraguay
Estilo: Paper-cut / Geometric con profundidad
Enfoque: MVP Motor de Cotización + Roadmap de productos
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pathlib import Path
import copy

# =============================================================================
# PALETA DE COLORES - Más vibrante y moderna
# =============================================================================
COLORS = {
    'primary_dark': RGBColor(15, 23, 42),      # Slate 900
    'primary': RGBColor(30, 64, 175),          # Blue 800
    'primary_light': RGBColor(59, 130, 246),   # Blue 500
    'accent': RGBColor(16, 185, 129),          # Emerald 500
    'accent_light': RGBColor(167, 243, 208),   # Emerald 200
    'warning': RGBColor(245, 158, 11),         # Amber 500
    'white': RGBColor(255, 255, 255),
    'gray_50': RGBColor(248, 250, 252),        # Slate 50
    'gray_100': RGBColor(241, 245, 249),       # Slate 100
    'gray_300': RGBColor(203, 213, 225),       # Slate 300
    'gray_600': RGBColor(71, 85, 105),         # Slate 600
    'gray_800': RGBColor(30, 41, 59),          # Slate 800
    'paper_cream': RGBColor(254, 252, 247),    # Cream for paper effect
    'paper_shadow': RGBColor(226, 232, 240),   # Shadow color
    'highlight': RGBColor(251, 191, 36),       # Amber 400
}

def add_shadow_shape(slide, shape, offset=Inches(0.05)):
    """Simula sombra agregando una forma detrás"""
    shadow = slide.shapes.add_shape(
        shape.shape_type,
        shape.left + offset,
        shape.top + offset,
        shape.width,
        shape.height
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = COLORS['paper_shadow']
    shadow.line.fill.background()
    # Mover la sombra detrás
    spTree = slide.shapes._spTree
    sp = shadow._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shadow

def create_paper_card(slide, x, y, width, height, color=None):
    """Crear una tarjeta con efecto de papel cortado"""
    if color is None:
        color = COLORS['white']
    
    # Sombra
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x + Inches(0.08), y + Inches(0.08),
        width, height
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = COLORS['paper_shadow']
    shadow.line.fill.background()
    
    # Tarjeta principal
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    
    # Ajustar radio de esquinas
    try:
        card.adjustments[0] = 0.05
        shadow.adjustments[0] = 0.05
    except:
        pass
    
    return card

def create_geometric_accent(slide, x, y, size, color, shape_type=MSO_SHAPE.OVAL):
    """Crear acento geométrico decorativo"""
    shape = slide.shapes.add_shape(shape_type, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, x, y, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    """Helper para agregar texto"""
    box = slide.shapes.add_textbox(x, y, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box

# =============================================================================
# SLIDES
# =============================================================================

def slide_cover(prs):
    """Slide 1: Portada con diseño geométrico"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo degradado simulado con rectángulos
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary_dark']
    bg.line.fill.background()
    
    # Formas geométricas decorativas (efecto paper-cut)
    # Círculo grande semi-transparente
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-1), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = COLORS['primary']
    circle1.line.fill.background()
    
    # Círculo mediano
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4), Inches(3), Inches(3))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['primary_light']
    circle2.line.fill.background()
    
    # Triángulo decorativo
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(8.5), Inches(5.5), Inches(2), Inches(2))
    tri.fill.solid()
    tri.fill.fore_color.rgb = COLORS['accent']
    tri.line.fill.background()
    
    # Rectángulo accent
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(4), Inches(0.7))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLORS['accent']
    rect.line.fill.background()
    
    # Título principal
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(8), Inches(1.5),
                 "Seguros Inteligentes", 48, COLORS['white'], bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(8), Inches(1),
                 "para HomeBanking", 44, COLORS['accent_light'], bold=True)
    
    # Subtítulo
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(7), Inches(0.8),
                 "El primer motor de cotización conversacional con IA\npara el mercado paraguayo", 
                 20, COLORS['gray_300'])
    
    # Badge MVP
    badge = create_paper_card(slide, Inches(0.8), Inches(5.8), Inches(2.5), Inches(0.6), COLORS['warning'])
    add_text_box(slide, Inches(0.9), Inches(5.9), Inches(2.3), Inches(0.4),
                 "MVP LISTO PARA DEMO", 14, COLORS['primary_dark'], bold=True, align=PP_ALIGN.CENTER)

def slide_opportunity(prs):
    """Slide 2: La Oportunidad en Paraguay"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['gray_50']
    bg.line.fill.background()
    
    # Header con forma geométrica
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Triángulo decorativo
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8), Inches(1.3), Inches(2), Inches(1.5))
    tri.fill.solid()
    tri.fill.fore_color.rgb = COLORS['accent']
    tri.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8), Inches(1),
                 "La Oportunidad en Paraguay", 32, COLORS['white'], bold=True)
    
    # Card principal - El contexto
    create_paper_card(slide, Inches(0.5), Inches(2.2), Inches(9), Inches(1.8), COLORS['white'])
    
    add_text_box(slide, Inches(0.8), Inches(2.4), Inches(8.5), Inches(0.5),
                 "El Momento es Ahora", 22, COLORS['primary'], bold=True)
    
    context_text = """• Paraguay no permite pólizas online actualmente
• Su banco está gestionando ser EL PRIMERO en obtener el permiso
• Ya tienen sistema interno para agentes en sucursales
• Falta el puente: cotización digital inteligente para clientes"""
    
    add_text_box(slide, Inches(0.8), Inches(2.95), Inches(8.5), Inches(1.4),
                 context_text, 16, COLORS['gray_600'])
    
    # Tres cards de estadísticas
    stats = [
        ("88%", "de aseguradoras\nglobales usan AI", "McKinsey"),
        ("58%", "de clientes quieren\nasesoría personalizada", "Accenture"),
        ("0", "bancos en Paraguay\ncon seguros digitales", "Oportunidad"),
    ]
    
    for i, (value, label, source) in enumerate(stats):
        x = Inches(0.5) + i * Inches(3.15)
        create_paper_card(slide, x, Inches(4.3), Inches(3), Inches(2.4), COLORS['white'])
        
        # Círculo con valor
        circle_color = [COLORS['primary_light'], COLORS['accent'], COLORS['warning']][i]
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(4.5), Inches(1.2), Inches(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = circle_color
        circle.line.fill.background()
        
        add_text_box(slide, x + Inches(0.9), Inches(4.75), Inches(1.2), Inches(0.7),
                     value, 28, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
        
        add_text_box(slide, x + Inches(0.2), Inches(5.8), Inches(2.6), Inches(0.8),
                     label, 13, COLORS['gray_600'], align=PP_ALIGN.CENTER)
        
        add_text_box(slide, x + Inches(0.2), Inches(6.4), Inches(2.6), Inches(0.3),
                     f"Fuente: {source}", 10, COLORS['gray_300'], align=PP_ALIGN.CENTER)

def slide_mvp_hero(prs):
    """Slide 3: MVP - Motor de Cotización (HERO slide)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo con acento
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # Forma decorativa grande
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLORS['primary_dark']
    deco.line.fill.background()
    
    # Card principal blanca
    create_paper_card(slide, Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.5), COLORS['white'])
    
    # Badge "LISTO PARA DEMO"
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.7), Inches(2.7), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS['accent']
    badge.line.fill.background()
    add_text_box(slide, Inches(6.6), Inches(1.78), Inches(2.5), Inches(0.4),
                 "LISTO PARA DEMO", 12, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
    
    # Título
    add_text_box(slide, Inches(1), Inches(1.8), Inches(5), Inches(0.6),
                 "Motor de Cotización", 28, COLORS['primary_dark'], bold=True)
    add_text_box(slide, Inches(1), Inches(2.35), Inches(5), Inches(0.5),
                 "Conversacional con IA", 24, COLORS['primary_light'], bold=True)
    
    # Descripción
    desc = """Tu cliente abre el app, escribe "quiero asegurar mi auto" 
y en 3 minutos tiene una cotización personalizada.

Sin formularios. Sin esperas. Sin fricciones."""
    add_text_box(slide, Inches(1), Inches(3.1), Inches(4.5), Inches(1.5),
                 desc, 15, COLORS['gray_600'])
    
    # Features en mini-cards
    features = [
        ("Conversación Natural", "Chat inteligente que guía al cliente"),
        ("Multi-Producto", "Auto, vida, hogar, salud - un solo motor"),
        ("Multi-LLM", "OpenAI, Anthropic, o modelos locales"),
        ("API Lista", "Integración simple con su sistema actual"),
    ]
    
    for i, (title, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = Inches(1) + col * Inches(4)
        y = Inches(4.7) + row * Inches(1.1)
        
        # Mini círculo
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.35), Inches(0.35))
        c.fill.solid()
        c.fill.fore_color.rgb = COLORS['accent']
        c.line.fill.background()
        
        add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(0.35), Inches(0.3),
                     "✓", 14, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
        
        add_text_box(slide, x + Inches(0.5), y, Inches(3.3), Inches(0.4),
                     title, 14, COLORS['primary_dark'], bold=True)
        add_text_box(slide, x + Inches(0.5), y + Inches(0.4), Inches(3.3), Inches(0.4),
                     desc, 11, COLORS['gray_600'])
    
    # Tiempo a producción
    time_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(3), Inches(3.3), Inches(1.4))
    time_box.fill.solid()
    time_box.fill.fore_color.rgb = COLORS['gray_100']
    time_box.line.fill.background()
    
    add_text_box(slide, Inches(6), Inches(3.15), Inches(3), Inches(0.4),
                 "Tiempo a Producción", 12, COLORS['gray_600'], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6), Inches(3.5), Inches(3), Inches(0.6),
                 "2-4 semanas", 26, COLORS['accent'], bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6), Inches(4), Inches(3), Inches(0.3),
                 "MVP ~80% completado", 11, COLORS['gray_600'], align=PP_ALIGN.CENTER)

def slide_how_it_works(prs):
    """Slide 4: Cómo Funciona"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['gray_50']
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.8),
                 "¿Cómo Funciona?", 28, COLORS['white'], bold=True)
    
    # Flujo con flechas
    steps = [
        ("1", "Cliente abre\nel app", COLORS['primary_light']),
        ("2", "Escribe:\n\"Quiero seguro\"", COLORS['primary']),
        ("3", "AI hace\npreguntas", COLORS['accent']),
        ("4", "Genera\ncotización", COLORS['warning']),
        ("5", "Cliente\nelige plan", COLORS['accent']),
    ]
    
    for i, (num, text, color) in enumerate(steps):
        x = Inches(0.3) + i * Inches(1.95)
        
        # Card
        create_paper_card(slide, x, Inches(1.9), Inches(1.8), Inches(2.2), COLORS['white'])
        
        # Número en círculo
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), Inches(2.1), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        add_text_box(slide, x + Inches(0.6), Inches(2.2), Inches(0.6), Inches(0.5),
                     num, 20, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
        
        add_text_box(slide, x + Inches(0.1), Inches(2.85), Inches(1.6), Inches(1),
                     text, 13, COLORS['gray_600'], align=PP_ALIGN.CENTER)
        
        # Flecha (excepto último)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.85), Inches(2.85), Inches(0.4), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['gray_300']
            arrow.line.fill.background()
    
    # Ejemplo de conversación
    create_paper_card(slide, Inches(0.5), Inches(4.4), Inches(9), Inches(2.8), COLORS['white'])
    
    add_text_box(slide, Inches(0.8), Inches(4.6), Inches(4), Inches(0.4),
                 "Ejemplo de Conversación", 16, COLORS['primary_dark'], bold=True)
    
    # Burbujas de chat
    # Usuario
    user_bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5), Inches(4.3), Inches(0.5))
    user_bubble.fill.solid()
    user_bubble.fill.fore_color.rgb = COLORS['primary_light']
    user_bubble.line.fill.background()
    add_text_box(slide, Inches(4.1), Inches(5.08), Inches(4.1), Inches(0.4),
                 "Quiero asegurar mi Honda CR-V 2024", 12, COLORS['white'])
    
    # AI
    ai_bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(5), Inches(0.5))
    ai_bubble.fill.solid()
    ai_bubble.fill.fore_color.rgb = COLORS['gray_100']
    ai_bubble.line.fill.background()
    add_text_box(slide, Inches(0.9), Inches(5.68), Inches(4.8), Inches(0.4),
                 "¡Excelente elección! ¿Lo usarás para ir al trabajo?", 12, COLORS['gray_800'])
    
    # Usuario
    user_bubble2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(6.2), Inches(2.8), Inches(0.5))
    user_bubble2.fill.solid()
    user_bubble2.fill.fore_color.rgb = COLORS['primary_light']
    user_bubble2.line.fill.background()
    add_text_box(slide, Inches(5.6), Inches(6.28), Inches(2.6), Inches(0.4),
                 "Sí, manejo 30km al día", 12, COLORS['white'])
    
    # Indicador
    add_text_box(slide, Inches(0.8), Inches(6.85), Inches(3), Inches(0.3),
                 "→ 3 minutos hasta cotización", 11, COLORS['accent'], bold=True)

def slide_vs_actual(prs):
    """Slide 5: Su Sistema Actual vs. Nuestra Solución"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['gray_50']
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.8),
                 "Su Sistema Actual + Nuestra Solución", 26, COLORS['white'], bold=True)
    
    # Card izquierda - Actual
    create_paper_card(slide, Inches(0.4), Inches(1.7), Inches(4.5), Inches(5.3), COLORS['white'])
    
    # Header card
    actual_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.7), Inches(4.5), Inches(0.7))
    actual_header.fill.solid()
    actual_header.fill.fore_color.rgb = COLORS['gray_600']
    actual_header.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(1.85), Inches(4), Inches(0.5),
                 "HOY: Sistema en Sucursales", 16, COLORS['white'], bold=True)
    
    actual_items = [
        "✓ Agentes capacitados",
        "✓ Pueden crear pólizas",
        "✓ Sistema interno funcionando",
        "✗ Solo presencial",
        "✗ Sin canal digital",
        "✗ Sin autoservicio",
        "✗ Horario limitado",
    ]
    
    for i, item in enumerate(actual_items):
        color = COLORS['accent'] if item.startswith("✓") else COLORS['gray_300']
        add_text_box(slide, Inches(0.7), Inches(2.6 + i * 0.55), Inches(4), Inches(0.5),
                     item, 14, color if item.startswith("✓") else COLORS['gray_600'])
    
    # Flecha central
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.7), Inches(4), Inches(0.6), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS['accent']
    arrow.line.fill.background()
    
    # Card derecha - Con nosotros
    create_paper_card(slide, Inches(5.1), Inches(1.7), Inches(4.5), Inches(5.3), COLORS['white'])
    
    # Header card
    new_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(1.7), Inches(4.5), Inches(0.7))
    new_header.fill.solid()
    new_header.fill.fore_color.rgb = COLORS['accent']
    new_header.line.fill.background()
    
    add_text_box(slide, Inches(5.3), Inches(1.85), Inches(4), Inches(0.5),
                 "MAÑANA: + Motor de Cotización", 15, COLORS['white'], bold=True)
    
    new_items = [
        "✓ Todo lo anterior PLUS:",
        "✓ Cotización en el app 24/7",
        "✓ Conversación natural con AI",
        "✓ Pre-calificación instantánea",
        "✓ Cliente llega a sucursal listo",
        "✓ Agente cierra más rápido",
        "✓ Datos ya capturados",
    ]
    
    for i, item in enumerate(new_items):
        weight = True if i == 0 else False
        add_text_box(slide, Inches(5.4), Inches(2.6 + i * 0.55), Inches(4), Inches(0.5),
                     item, 14, COLORS['accent'] if i > 0 else COLORS['primary_dark'], bold=weight)

def slide_roadmap_productos(prs):
    """Slide 6: Roadmap de Productos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary_dark']
    bg.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.8),
                 "Roadmap de Desarrollo", 28, COLORS['white'], bold=True)
    
    add_text_box(slide, Inches(0.6), Inches(0.95), Inches(8), Inches(0.4),
                 "Tiempos estimados con 1 desarrollador dedicado", 14, COLORS['gray_300'])
    
    # Timeline
    products = [
        ("Motor de Cotización", "2-4 sem", "MVP 80% listo", COLORS['accent'], True),
        ("Document Intelligence", "6-8 sem", "Extracción de docs", COLORS['primary_light'], False),
        ("Analizador de Productos", "4-6 sem", "Comparador inteligente", COLORS['primary_light'], False),
        ("Copiloto para Agentes", "8-10 sem", "Asistente en sucursal", COLORS['primary_light'], False),
        ("Agente de Claims", "10-12 sem", "Procesamiento siniestros", COLORS['primary_light'], False),
        ("Dashboard Gobernanza", "6-8 sem", "Auditoría y métricas", COLORS['primary_light'], False),
    ]
    
    for i, (name, time, desc, color, is_mvp) in enumerate(products):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * Inches(4.8)
        y = Inches(1.6) + row * Inches(1.8)
        
        # Card
        card_color = COLORS['accent'] if is_mvp else COLORS['gray_800']
        create_paper_card(slide, x, y, Inches(4.5), Inches(1.5), card_color if is_mvp else COLORS['white'])
        
        if is_mvp:
            # Badge NOW
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(3.4), y + Inches(0.15), Inches(0.9), Inches(0.35))
            badge.fill.solid()
            badge.fill.fore_color.rgb = COLORS['warning']
            badge.line.fill.background()
            add_text_box(slide, x + Inches(3.45), y + Inches(0.2), Inches(0.8), Inches(0.3),
                         "AHORA", 10, COLORS['primary_dark'], bold=True, align=PP_ALIGN.CENTER)
            
            text_color = COLORS['white']
        else:
            text_color = COLORS['gray_800']
        
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3), Inches(0.4),
                     name, 15, text_color, bold=True)
        
        # Tiempo en círculo
        time_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.7), Inches(1.3), Inches(0.5))
        time_bg.fill.solid()
        time_bg.fill.fore_color.rgb = COLORS['primary_dark'] if is_mvp else COLORS['gray_100']
        time_bg.line.fill.background()
        
        add_text_box(slide, x + Inches(0.25), y + Inches(0.78), Inches(1.2), Inches(0.4),
                     time, 12, COLORS['white'] if is_mvp else COLORS['gray_600'], bold=True, align=PP_ALIGN.CENTER)
        
        add_text_box(slide, x + Inches(1.6), y + Inches(0.8), Inches(2.7), Inches(0.4),
                     desc, 12, text_color if is_mvp else COLORS['gray_600'])

def slide_business_model(prs):
    """Slide 7: Modelo de Negocio - 3 Opciones"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['gray_50']
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.8),
                 "Modelo de Negocio: 3 Opciones", 28, COLORS['white'], bold=True)
    
    # Tres opciones
    options = [
        ("Opción A", "Bajo Riesgo", [
            "100% desarrollo pagado",
            "0.5-1% comisión por póliza",
            "Soporte: cobro por hora",
            "Ideal: control total de costos"
        ], COLORS['primary_light']),
        ("Opción B", "Balanced", [
            "50% desarrollo pagado",
            "1.5-2% comisión por póliza",
            "Soporte básico incluido",
            "Ideal: riesgo compartido"
        ], COLORS['accent']),
        ("Opción C", "Partnership", [
            "$0 desarrollo inicial",
            "3-5% comisión por póliza",
            "Todo incluido",
            "Ideal: sin inversión upfront"
        ], COLORS['warning']),
    ]
    
    for i, (title, subtitle, items, color) in enumerate(options):
        x = Inches(0.4) + i * Inches(3.2)
        
        # Card
        create_paper_card(slide, x, Inches(1.7), Inches(3), Inches(4.8), COLORS['white'])
        
        # Header de card
        card_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.7), Inches(3), Inches(1))
        card_header.fill.solid()
        card_header.fill.fore_color.rgb = color
        card_header.line.fill.background()
        
        add_text_box(slide, x + Inches(0.1), Inches(1.85), Inches(2.8), Inches(0.4),
                     title, 18, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(2.25), Inches(2.8), Inches(0.35),
                     subtitle, 13, COLORS['white'], align=PP_ALIGN.CENTER)
        
        # Items
        for j, item in enumerate(items):
            is_highlight = "%" in item or "$" in item
            add_text_box(slide, x + Inches(0.2), Inches(2.9 + j * 0.65), Inches(2.6), Inches(0.6),
                         "• " + item, 12, COLORS['primary_dark'] if is_highlight else COLORS['gray_600'],
                         bold=is_highlight)
        
        # Recomendado badge para opción B
        if i == 1:
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.5), Inches(5.9), Inches(2), Inches(0.45))
            badge.fill.solid()
            badge.fill.fore_color.rgb = COLORS['accent']
            badge.line.fill.background()
            add_text_box(slide, x + Inches(0.6), Inches(5.98), Inches(1.8), Inches(0.35),
                         "RECOMENDADO", 11, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
    
    # Nota al pie
    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(9), Inches(0.4),
                 "Todas las opciones incluyen: documentación, capacitación inicial, y 30 días de soporte post-launch", 
                 11, COLORS['gray_600'])

def slide_roi(prs):
    """Slide 8: Proyección de ROI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['gray_50']
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.8),
                 "¿Por qué Invertir? El ROI", 28, COLORS['white'], bold=True)
    
    # Dos columnas
    # Izquierda - Reducción de costos
    create_paper_card(slide, Inches(0.4), Inches(1.7), Inches(4.5), Inches(3.2), COLORS['white'])
    
    cost_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.7), Inches(4.5), Inches(0.7))
    cost_header.fill.solid()
    cost_header.fill.fore_color.rgb = COLORS['primary_light']
    cost_header.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(1.85), Inches(4), Inches(0.5),
                 "Reducción de Costos", 16, COLORS['white'], bold=True)
    
    costs = [
        "-80% costo por cotización",
        "Menos tiempo de agente por cliente",
        "Datos pre-capturados = menos errores",
        "24/7 sin costo de personal adicional",
    ]
    for i, item in enumerate(costs):
        add_text_box(slide, Inches(0.6), Inches(2.6 + i * 0.5), Inches(4.1), Inches(0.45),
                     "• " + item, 13, COLORS['gray_600'])
    
    # Derecha - Aumento de ingresos
    create_paper_card(slide, Inches(5.1), Inches(1.7), Inches(4.5), Inches(3.2), COLORS['white'])
    
    rev_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(1.7), Inches(4.5), Inches(0.7))
    rev_header.fill.solid()
    rev_header.fill.fore_color.rgb = COLORS['accent']
    rev_header.line.fill.background()
    
    add_text_box(slide, Inches(5.3), Inches(1.85), Inches(4), Inches(0.5),
                 "Aumento de Ingresos", 16, COLORS['white'], bold=True)
    
    revenues = [
        "Nuevo revenue stream (comisiones)",
        "Cross-sell en app: crédito → seguro",
        "Mayor conversión por experiencia",
        "Primeros en el mercado = ventaja",
    ]
    for i, item in enumerate(revenues):
        add_text_box(slide, Inches(5.3), Inches(2.6 + i * 0.5), Inches(4.1), Inches(0.45),
                     "• " + item, 13, COLORS['gray_600'])
    
    # Stats de referencia
    create_paper_card(slide, Inches(0.4), Inches(5.1), Inches(9.2), Inches(1.8), COLORS['primary_dark'])
    
    add_text_box(slide, Inches(0.7), Inches(5.3), Inches(8.5), Inches(0.4),
                 "Resultados demostrados en la industria:", 14, COLORS['white'], bold=True)
    
    ref_stats = [
        (">30%", "eficiencia en underwriting", "BCG"),
        ("91%", "claims automatizados", "Dearborn Labs"),
        ("3X", "más eficiencia operativa", "Dearborn Labs"),
    ]
    
    for i, (val, label, source) in enumerate(ref_stats):
        x = Inches(0.7) + i * Inches(3.1)
        add_text_box(slide, x, Inches(5.75), Inches(1), Inches(0.5),
                     val, 22, COLORS['accent'], bold=True)
        add_text_box(slide, x + Inches(1.1), Inches(5.8), Inches(2), Inches(0.7),
                     f"{label}\n({source})", 11, COLORS['gray_300'])

def slide_why_us(prs):
    """Slide 9: ¿Por qué Nosotros?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # Decoración
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(-1), Inches(4), Inches(4))
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLORS['primary_dark']
    deco.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8), Inches(0.8),
                 "¿Por qué Nosotros?", 32, COLORS['white'], bold=True)
    
    # Cards de diferenciadores
    diffs = [
        ("MVP Funcionando", "No es PowerPoint.\nTenemos código en producción.", "✓"),
        ("Enfoque Paraguay", "Entendemos el contexto local\ny la regulación en proceso.", "🇵🇾"),
        ("Tecnología Flexible", "Multi-LLM, APIs abiertas,\nintegración con su sistema.", "⚙"),
        ("Partnership Real", "No somos vendor.\nCrecemos juntos.", "🤝"),
    ]
    
    for i, (title, desc, icon) in enumerate(diffs):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * Inches(4.8)
        y = Inches(1.5) + row * Inches(2.5)
        
        create_paper_card(slide, x, y, Inches(4.5), Inches(2.2), COLORS['white'])
        
        # Icono
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.3), Inches(0.7), Inches(0.7))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = COLORS['accent']
        icon_bg.line.fill.background()
        
        add_text_box(slide, x + Inches(0.25), y + Inches(0.4), Inches(0.6), Inches(0.5),
                     icon, 20, COLORS['white'], align=PP_ALIGN.CENTER)
        
        add_text_box(slide, x + Inches(1.1), y + Inches(0.35), Inches(3.2), Inches(0.5),
                     title, 16, COLORS['primary_dark'], bold=True)
        
        add_text_box(slide, x + Inches(1.1), y + Inches(0.85), Inches(3.2), Inches(1.1),
                     desc, 13, COLORS['gray_600'])

def slide_next_steps(prs):
    """Slide 10: Próximos Pasos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['gray_50']
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.8),
                 "Próximos Pasos", 28, COLORS['white'], bold=True)
    
    # Timeline
    steps = [
        ("ESTA SEMANA", "Demo en vivo del MVP", "Les mostramos el motor funcionando con sus productos"),
        ("SEMANA 2", "Definición de alcance", "Elegimos modelo de negocio y productos prioritarios"),
        ("SEMANA 3-4", "Integración técnica", "Conectamos con su sistema actual"),
        ("SEMANA 5-6", "Piloto controlado", "Prueba con grupo de clientes seleccionados"),
        ("SEMANA 7+", "Lanzamiento", "Listos para cuando llegue la autorización"),
    ]
    
    for i, (time, title, desc) in enumerate(steps):
        y = Inches(1.7) + i * Inches(1.05)
        
        # Línea de tiempo
        if i < len(steps) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.45), y + Inches(0.7), Inches(0.04), Inches(0.6))
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS['gray_300']
            line.line.fill.background()
        
        # Círculo
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.15), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['accent'] if i == 0 else COLORS['primary_light']
        circle.line.fill.background()
        
        add_text_box(slide, Inches(1.15), y + Inches(0.28), Inches(0.6), Inches(0.5),
                     str(i + 1), 18, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
        
        # Card
        create_paper_card(slide, Inches(2), y, Inches(7.5), Inches(0.9), COLORS['white'])
        
        add_text_box(slide, Inches(2.2), y + Inches(0.08), Inches(1.8), Inches(0.35),
                     time, 11, COLORS['accent'], bold=True)
        
        add_text_box(slide, Inches(4), y + Inches(0.08), Inches(5.3), Inches(0.35),
                     title, 14, COLORS['primary_dark'], bold=True)
        
        add_text_box(slide, Inches(4), y + Inches(0.45), Inches(5.3), Inches(0.4),
                     desc, 11, COLORS['gray_600'])

def slide_cta(prs):
    """Slide 11: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary_dark']
    bg.line.fill.background()
    
    # Formas decorativas
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = COLORS['primary']
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-2), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['accent']
    circle2.line.fill.background()
    
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(8), Inches(5), Inches(2.5), Inches(2.5))
    tri.fill.solid()
    tri.fill.fore_color.rgb = COLORS['warning']
    tri.line.fill.background()
    
    # Texto principal
    add_text_box(slide, Inches(0.5), Inches(2), Inches(9), Inches(1.2),
                 "¿Listos para ser los primeros\nen seguros digitales en Paraguay?", 
                 36, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
    
    # CTA
    cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(4), Inches(4), Inches(0.8))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = COLORS['accent']
    cta_box.line.fill.background()
    
    add_text_box(slide, Inches(3), Inches(4.15), Inches(4), Inches(0.6),
                 "Agendemos la Demo →", 22, COLORS['white'], bold=True, align=PP_ALIGN.CENTER)
    
    # Contacto
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.5),
                 "El MVP está listo. El momento es ahora.", 
                 18, COLORS['gray_300'], align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
                 "contacto@insurtech.com.py", 
                 16, COLORS['accent_light'], align=PP_ALIGN.CENTER)

# =============================================================================
# MAIN
# =============================================================================

def create_presentation():
    """Crear la presentación completa"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Generar slides
    slide_cover(prs)           # 1. Portada
    slide_opportunity(prs)     # 2. La Oportunidad en Paraguay
    slide_mvp_hero(prs)        # 3. MVP Hero
    slide_how_it_works(prs)    # 4. Cómo Funciona
    slide_vs_actual(prs)       # 5. Actual vs. Nuevo
    slide_roadmap_productos(prs)  # 6. Roadmap
    slide_business_model(prs)  # 7. Modelo de Negocio
    slide_roi(prs)             # 8. ROI
    slide_why_us(prs)          # 9. Por qué nosotros
    slide_next_steps(prs)      # 10. Próximos pasos
    slide_cta(prs)             # 11. CTA
    
    # Guardar
    output_path = Path(__file__).parent / "pitch-seguros-homebanking-v2.pptx"
    prs.save(output_path)
    print(f"Presentación generada: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
